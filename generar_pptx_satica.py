# =============================================================================
#  SATICA V2.0 — GENERADOR DE PRESENTACIÓN PREMIUM (python-pptx)
#  Fondo oscuro real · Tipografía moderna · Gráficos de barras integrados
# =============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, math, sys

# ── Paleta ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0f, 0x17, 0x2a)   # fondo marino oscuro
CARD      = RGBColor(0x1e, 0x29, 0x3b)   # tarjetas
ACCENT    = RGBColor(0x38, 0xbd, 0xf8)   # azul cielo
WHITE     = RGBColor(0xf8, 0xfa, 0xfc)
SUBTITLE  = RGBColor(0x94, 0xa3, 0xb8)
RED       = RGBColor(0xef, 0x44, 0x44)
ORANGE    = RGBColor(0xf5, 0x9e, 0x0b)
YELLOW    = RGBColor(0xfa, 0xcc, 0x15)
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
PURPLE    = RGBColor(0xa8, 0x55, 0xf7)
DARK_LINE = RGBColor(0x33, 0x41, 0x55)
FOOTER_C  = RGBColor(0x47, 0x55, 0x69)

# ── Dimensiones (widescreen 16:9) ────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout vacío

# =============================================================================
#  HELPERS
# =============================================================================

def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_slide_bg(slide, color: RGBColor):
    """Pinta el fondo completo de la diapositiva."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Agrega un rectángulo con relleno y/o borde."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=14, color=WHITE, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri",
                 line_spacing=None):
    """Cuadro de texto simple."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', f'{int(line_spacing*100000)}')
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txBox

def add_multiline_box(slide, lines, l, t, w, h,
                      font_size=12, color=WHITE, bold=False,
                      align=PP_ALIGN.LEFT, font_name="Calibri",
                      line_space_pct=130):
    """Cuadro con múltiples párrafos."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, clr, sz, bd) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        # Espaciado entre líneas
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(line_space_pct * 1000))
        run = p.add_run()
        run.text = txt
        run.font.name  = font_name
        run.font.size  = Pt(sz)
        run.font.color.rgb = clr
        run.font.bold  = bd
    return txBox

def gradient_rect(slide, l, t, w, h, color1, color2, angle=0):
    """Rectángulo con gradiente lineal."""
    shape = add_rect(slide, l, t, w, h)
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = color2
    return shape

def bar_chart_shape(slide, data, l, t, w, h,
                    bar_colors=None, show_values=True,
                    base_line=None, base_label=None):
    """
    Dibuja gráfica de barras verticales como formas nativas PowerPoint.
    data: list of (label, value)  — value entre 0 y 100
    """
    n      = len(data)
    max_v  = 100
    pad_l  = Inches(0.4)
    pad_r  = Inches(0.15)
    pad_t  = Inches(0.25)
    pad_b  = Inches(0.55)
    chart_w = w - pad_l - pad_r
    chart_h = h - pad_t - pad_b
    bar_w   = chart_w / (n * 1.65)
    gap     = (chart_w - bar_w * n) / (n + 1)

    # Ejes (líneas base)
    add_rect(slide, l + pad_l, t + pad_t,
             chart_w, Pt(1.5), DARK_LINE)   # línea tope
    add_rect(slide, l + pad_l, t + pad_t + chart_h,
             chart_w, Pt(2), DARK_LINE)      # línea base

    if base_line is not None:
        by = t + pad_t + chart_h * (1 - base_line / max_v)
        add_rect(slide, l + pad_l, by, chart_w, Pt(1.5),
                 RGBColor(0x64, 0x74, 0x8b))
        if base_label:
            add_text_box(slide, base_label,
                         l + pad_l, by - Inches(0.22),
                         Inches(1.3), Inches(0.2),
                         font_size=7, color=RGBColor(0x64,0x74,0x8b))

    for i, (label, value) in enumerate(data):
        bx = l + pad_l + gap + i * (bar_w + gap)
        bh = chart_h * (value / max_v)
        by = t + pad_t + chart_h - bh
        col = bar_colors[i] if bar_colors else ACCENT

        # Sombra suave
        add_rect(slide, bx + Inches(0.03), by + Inches(0.04),
                 bar_w, bh, RGBColor(0,0,0))

        # Barra
        shape = add_rect(slide, bx, by, bar_w, bh, col)

        # Valor encima
        if show_values:
            add_text_box(slide, f"{value}%",
                         bx - Inches(0.05), by - Inches(0.28),
                         bar_w + Inches(0.1), Inches(0.25),
                         font_size=9.5, color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER)

        # Etiqueta abajo
        add_text_box(slide, label,
                     bx - Inches(0.08), t + pad_t + chart_h + Inches(0.05),
                     bar_w + Inches(0.16), pad_b - Inches(0.05),
                     font_size=7.5, color=SUBTITLE, bold=False,
                     align=PP_ALIGN.CENTER)

def line_chart(slide, data, l, t, w, h, dot_colors=None):
    """Gráfica de líneas simples para recall acumulado."""
    n       = len(data)
    pad_l, pad_r, pad_t, pad_b = Inches(0.35), Inches(0.15), Inches(0.3), Inches(0.55)
    chart_w = w - pad_l - pad_r
    chart_h = h - pad_t - pad_b
    max_v   = 100
    step_x  = chart_w / (n - 1)

    pts = []
    for i, (label, value) in enumerate(data):
        px = l + pad_l + i * step_x
        py = t + pad_t + chart_h * (1 - value / max_v)
        pts.append((px, py, label, value))

    # Líneas entre puntos
    for i in range(len(pts) - 1):
        x1, y1 = pts[i][0], pts[i][1]
        x2, y2 = pts[i+1][0], pts[i+1][1]
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(3.5)

    # Puntos y etiquetas
    for i, (px, py, label, value) in enumerate(pts):
        col = dot_colors[i] if dot_colors else ACCENT
        dot_r = Inches(0.13)
        # Sombra
        add_rect(slide, px - dot_r + Inches(0.03),
                 py - dot_r + Inches(0.04),
                 dot_r*2, dot_r*2, RGBColor(0,0,0))
        # Punto relleno
        sh = slide.shapes.add_shape(9, px-dot_r, py-dot_r, dot_r*2, dot_r*2)
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        sh.line.fill.background()
        # Valor
        add_text_box(slide, f"{value}%",
                     px - Inches(0.22), py - Inches(0.35),
                     Inches(0.44), Inches(0.25),
                     font_size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Etiqueta
        add_text_box(slide, label,
                     px - Inches(0.45), t + pad_t + chart_h + Inches(0.05),
                     Inches(0.9), pad_b,
                     font_size=7.5, color=SUBTITLE, align=PP_ALIGN.CENTER)

def title_slide_header(slide, title, subtitle=None,
                       tc=ACCENT, sc=SUBTITLE, tsz=26):
    add_text_box(slide, title,
                 Inches(0.45), Inches(0.18), Inches(12.5), Inches(0.65),
                 font_size=tsz, color=tc, bold=True)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.45), Inches(0.82), Inches(12.5), Inches(0.35),
                     font_size=11.5, color=sc)
    # Línea divisoria
    add_rect(slide, Inches(0.45), Inches(1.13), Inches(12.4), Pt(1.6), DARK_LINE)

def accent_bar_left(slide, color=ACCENT):
    gradient_rect(slide, 0, 0, Inches(0.32), H, color, CARD, angle=270)

def footer_line(slide, text="CVC — DAR Suroriente  ·  SATICA V2.0  ·  Abril 2026"):
    add_rect(slide, 0, H - Inches(0.32), W, Inches(0.32), DARK_LINE)
    add_text_box(slide, text, Inches(0.45), H - Inches(0.28), Inches(12), Inches(0.25),
                 font_size=8, color=FOOTER_C, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h, fill=CARD, radius=None):
    sh = add_rect(slide, l, t, w, h, fill)
    return sh

# =============================================================================
#  SLIDE 1 — PORTADA
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)

# Gradiente lateral izquierdo
gradient_rect(sl, 0, 0, Inches(0.38), H, ACCENT, RGBColor(0x06,0x82,0xa8), angle=270)

# Banda superior sutil
gradient_rect(sl, Inches(0.38), 0, W, Inches(0.08),
              RGBColor(0x38,0xbd,0xf8), BG, angle=0)

# Organización
add_text_box(sl,
    "CORPORACIÓN AUTÓNOMA REGIONAL DEL VALLE DEL CAUCA",
    Inches(0.6), Inches(0.9), Inches(11.5), Inches(0.35),
    font_size=9, color=SUBTITLE, bold=True)

# Título gigante
add_text_box(sl, "SATICA",
    Inches(0.6), Inches(1.45), Inches(10), Inches(1.55),
    font_size=90, color=WHITE, bold=True)

# Versión badge
card(sl, Inches(5.8), Inches(1.5), Inches(1.2), Inches(0.5), ACCENT)
add_text_box(sl, "V 2.0",
    Inches(5.83), Inches(1.52), Inches(1.14), Inches(0.46),
    font_size=18, color=BG, bold=True, align=PP_ALIGN.CENTER)

# Subtítulo
add_text_box(sl,
    "Sistema de Alertas Tempranas de Incendios en Caña de Azúcar",
    Inches(0.6), Inches(3.1), Inches(10.5), Inches(0.55),
    font_size=22, color=ACCENT)

# Línea
add_rect(sl, Inches(0.6), Inches(3.72), Inches(9.5), Pt(2), ACCENT)

# Descripción
add_text_box(sl,
    "Monitoreo satelital continuo  ·  Inteligencia Artificial XGBoost V9  ·  Alertas en tiempo real",
    Inches(0.6), Inches(3.85), Inches(11), Inches(0.4),
    font_size=12.5, color=SUBTITLE)

# Stats portada
stats = [
    ("5", "Satélites", ACCENT),
    ("98.1%", "Certeza", RED),
    ("15 min", "Alertas", GREEN),
    ("8×", "vs Azar", PURPLE),
]
for i, (num, lbl, col) in enumerate(stats):
    sx = Inches(0.6) + i * Inches(2.45)
    card(sl, sx, Inches(4.55), Inches(2.25), Inches(1.42), CARD)
    add_text_box(sl, num,
        sx + Inches(0.08), Inches(4.65), Inches(2.1), Inches(0.75),
        font_size=34, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, lbl,
        sx + Inches(0.08), Inches(5.38), Inches(2.1), Inches(0.28),
        font_size=10.5, color=SUBTITLE, align=PP_ALIGN.CENTER)

# Pie
add_rect(sl, 0, H - Inches(0.36), W, Inches(0.36), DARK_LINE)
add_text_box(sl,
    "Abril 2026  ·  Presentación Ejecutiva  ·  DAR Suroriente",
    Inches(0.5), H - Inches(0.3), Inches(11), Inches(0.26),
    font_size=8.5, color=FOOTER_C)

print("  ✅ Slide 1 — Portada")

# =============================================================================
#  SLIDE 2 — ¿QUÉ ES SATICA?
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "¿QUÉ ES SATICA?",
    "El «Ojo de Halcón» de la CVC para proteger el medio ambiente del Valle del Cauca")

# Texto izquierda
card(sl, Inches(0.5), Inches(1.3), Inches(5.7), Inches(5.75), CARD)
add_multiline_box(sl, [
    ("SATICA es mucho más que un software.", ACCENT, 13.5, True),
    ("", WHITE, 4, False),
    ("Es una plataforma de inteligencia ambiental que fusiona datos de satélites NASA, modelos de IA y cartografía oficial de la CVC para predecir, detectar y responder a incendios en caña de azúcar antes de que se salgan de control.", SUBTITLE, 11.5, False),
    ("", WHITE, 6, False),
    ("MISIÓN AMBIENTAL", ACCENT, 11, True),
    ("", WHITE, 3, False),
    ("🔴  Detectar focos antes de que sean incontrolables.", WHITE, 11, False),
    ("🌫️  Reducir emisiones de material particulado PM2.5.", WHITE, 11, False),
    ("🦜  Proteger fauna, comunidades y ecosistemas.", WHITE, 11, False),
    ("⚖️  Blindaje jurídico para técnicos CVC (Res. 0741/2016).", WHITE, 11, False),
], Inches(0.65), Inches(1.42), Inches(5.4), Inches(5.5), line_space_pct=130)

# Stats derecha (2×3)
stats2 = [
    ("5",      "Satélites\nintegrados",         ACCENT,  36),
    ("98.1%",  "Certeza\nFactual",              RED,     28),
    ("6+",     "Años de datos\nhistóricos",     GREEN,   36),
    ("15 min", "Ciclo de\nalertas",             ORANGE,  24),
    ("8×",     "Más preciso\nque el azar",      PURPLE,  36),
    ("40%",    "Predios con\nriesgo silencioso",YELLOW,  28),
]
cols = [Inches(6.55), Inches(9.55)]
rows = [Inches(1.35), Inches(3.2), Inches(5.05)]
for idx, (num, lbl, col, fsz) in enumerate(stats2):
    ci, ri = idx % 2, idx // 2
    sx, sy = cols[ci], rows[ri]
    card(sl, sx, sy, Inches(2.6), Inches(1.6), CARD)
    add_text_box(sl, num, sx+Inches(0.1), sy+Inches(0.12),
                 Inches(2.4), Inches(0.85),
                 font_size=fsz, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, lbl, sx+Inches(0.1), sy+Inches(0.98),
                 Inches(2.4), Inches(0.55),
                 font_size=9.5, color=SUBTITLE, align=PP_ALIGN.CENTER)

footer_line(sl)
print("  ✅ Slide 2 — ¿Qué es SATICA?")

# =============================================================================
#  SLIDE 3 — FUNCIONES PRINCIPALES
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "FUNCIONES PRINCIPALES",
    "5 módulos integrados que conforman el sistema de inteligencia ambiental de la CVC")

funciones = [
    ("🛰️ RADAR\nSATELITAL",    RED,
     ["4 sensores NASA activos.", "Focos térmicos en < 15 min.",
      "Cruce automático contra\ncatastro CVC.", "Alertas Telegram + Deep Link."]),
    ("🌱 ANÁLISIS\nBIOMASA",    GREEN,
     ["Sentinel-2 a 10 m.", "NDVI = Humedad del cultivo.",
      "NBR = Cicatrices post-quema.", "Via Google Earth Engine."]),
    ("🧠 IA XGBOOST\nV9",       PURPLE,
     ["Entrenado 2019–2024.", "Opera con cielo nublado.",
      "Variables: vías, pueblos,\nrecurrencia.", "8× más preciso que el azar."]),
    ("💨 HUMO\nHYSPLIT",        ACCENT,
     ["Trayectoria a 6 horas.", "Datos viento Open-Meteo.",
      "Comunidades en riesgo.", "Visible en mapa interactivo."]),
    ("📋 CENTRO\nOPERACIONES",   ORANGE,
     ["Boletines visita PDF/CSV.", "Capas GIS KML/Shapefile.",
      "Red de riesgo por ingenio.", "Historial 2019–2026."]),
]

col_w = Inches(2.36)
for i, (icon, col, bullets) in enumerate(funciones):
    lx = Inches(0.48) + i * Inches(2.57)
    # Cabecera coloreada
    card(sl, lx, Inches(1.32), col_w, Inches(0.95), col)
    add_text_box(sl, icon, lx+Inches(0.05), Inches(1.35),
                 col_w-Inches(0.1), Inches(0.88),
                 font_size=12, color=BG, bold=True, align=PP_ALIGN.CENTER)
    # Cuerpo oscuro
    card(sl, lx, Inches(2.3), col_w, Inches(4.75), CARD)
    lines = []
    for b in bullets:
        lines.append(("• " + b, SUBTITLE, 10.5, False))
        lines.append(("", WHITE, 3, False))
    add_multiline_box(sl, lines,
                      lx+Inches(0.1), Inches(2.4),
                      col_w-Inches(0.2), Inches(4.55),
                      line_space_pct=125)

footer_line(sl)
print("  ✅ Slide 3 — Funciones")

# =============================================================================
#  SLIDE 4 — 5 SATÉLITES
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "5 SATÉLITES · 1 SISTEMA",
    "Cada sensor aporta una capa única de información. SATICA los funde en una sola alerta integrada.")

sats = [
    ("VIIRS\nSuomi-NPP",    RED,    "Focos térmicos activos",     "375 m resolución · NRT · Sensor MODIS mejorado"),
    ("VIIRS\nNOAA-20",      ORANGE, "Respaldo y redundancia",     "Segunda constelación VIIRS · Alta confianza dual"),
    ("MODIS\nTerra / Aqua", hex2rgb("#f97316"), "Potencia Radiativa del Fuego",  "FRP · 1 km · Intensidad y tamaño del incendio"),
    ("Sentinel-2\nCopernicus",ACCENT,"Biomasa y cicatrices",      "NDVI + NBR · 10 m · Google Earth Engine"),
    ("HYSPLIT\nNOAA",       PURPLE, "Trayectoria del humo",       "Predicción 6 horas · Open-Meteo · Sin clave API"),
]

row_h = Inches(1.04)
for i, (name, col, title, desc) in enumerate(sats):
    ty = Inches(1.3) + i * row_h
    # Badge de color
    card(sl, Inches(0.5), ty, Inches(1.95), row_h - Inches(0.08), col)
    add_text_box(sl, name, Inches(0.52), ty+Inches(0.06),
                 Inches(1.9), row_h-Inches(0.14),
                 font_size=11.5, color=BG, bold=True, align=PP_ALIGN.CENTER)
    # Título dato
    card(sl, Inches(2.5), ty, Inches(10.6), row_h - Inches(0.08), CARD)
    add_text_box(sl, title,
                 Inches(2.65), ty+Inches(0.04), Inches(4.0), Inches(0.42),
                 font_size=13, color=WHITE, bold=True)
    add_text_box(sl, desc,
                 Inches(2.65), ty+Inches(0.46), Inches(10.2), Inches(0.46),
                 font_size=10.5, color=SUBTITLE)
    # Número de sensor
    add_text_box(sl, str(i+1),
                 Inches(12.7), ty+Inches(0.12), Inches(0.5), Inches(0.7),
                 font_size=28, color=col, bold=True, align=PP_ALIGN.CENTER)

footer_line(sl)
print("  ✅ Slide 4 — Satélites")

# =============================================================================
#  SLIDE 5 — FLUJO: CÓMO FUNCIONA
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "¿CÓMO FUNCIONA SATICA?",
    "Desde el satélite hasta el celular del técnico en menos de 15 minutos")

pasos = [
    ("01", "Órbita\nSatelital",  RED,
     "VIIRS, MODIS y Sentinel-2 pasan sobre el Valle del Cauca y registran anomalías térmicas y de biomasa."),
    ("02", "NASA\nFIRMS API",    ORANGE,
     "SATICA consulta automáticamente cada 15 minutos la base de datos en vivo de la NASA."),
    ("03", "Cruce\nCatastral",   YELLOW,
     "Cada foco se cruza contra el mapa oficial de suertes de caña (SOR_OK.shp) de la CVC."),
    ("04", "Análisis\nde Riesgo",ACCENT,
     "El motor XGBoost V9 calcula la criticidad combinando historial + telemetría + distancias."),
    ("05", "Alerta\nTelegram",   GREEN,
     "Si el predio es CRÍTICO, se envía la Ficha de Acción Rápida con coordenadas exactas y deep link."),
]

card_w = Inches(2.36)
for i, (num, title, col, desc) in enumerate(pasos):
    lx = Inches(0.48) + i * Inches(2.57)
    # Círculo número
    sh = sl.shapes.add_shape(9, lx + Inches(0.58), Inches(1.3), Inches(1.2), Inches(1.2))
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()
    add_text_box(sl, num, lx+Inches(0.58), Inches(1.38), Inches(1.2), Inches(0.9),
                 font_size=28, color=BG, bold=True, align=PP_ALIGN.CENTER)
    # Flecha
    if i < len(pasos)-1:
        ax = lx + card_w + Inches(0.1)
        add_text_box(sl, "▶", ax, Inches(1.7), Inches(0.26), Inches(0.5),
                     font_size=14, color=DARK_LINE, align=PP_ALIGN.CENTER)
    # Título
    add_text_box(sl, title, lx, Inches(2.6), card_w, Inches(0.7),
                 font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Descripción
    card(sl, lx, Inches(3.38), card_w, Inches(3.7), CARD)
    add_text_box(sl, desc, lx+Inches(0.1), Inches(3.5),
                 card_w-Inches(0.2), Inches(3.5),
                 font_size=11, color=SUBTITLE, align=PP_ALIGN.CENTER, line_spacing=1.35)

footer_line(sl)
print("  ✅ Slide 5 — Flujo")

# =============================================================================
#  SLIDE 6 — CERTEZA CIENTÍFICA (gráfico de barras)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "CERTEZA CIENTÍFICA: 5 NIVELES DE EVIDENCIA",
    "Back-Testing riguroso: 6 años de entrenamiento (2019–2024) · 1 año de validación ciega (2025)",
    tc=RED, tsz=22)

# Gráfico de barras izquierda
bar_data = [
    ("N1\nFLAGRANCIA",  98.1),
    ("N2\nPATRÓN",      86.4),
    ("N3\nIA",          72.3),
    ("N4\nSEÑAL",       44.7),
    ("N5\nSIN DATOS",   11.9),
]
bar_cols = [RED, ORANGE, PURPLE, ACCENT, GREEN]
card(sl, Inches(0.48), Inches(1.28), Inches(6.7), Inches(5.8), CARD)
add_text_box(sl, "PRECISION: ¿Cuándo el sistema alerta, acierta?",
             Inches(0.62), Inches(1.35), Inches(6.4), Inches(0.28),
             font_size=9, color=FOOTER_C, bold=True)
bar_chart_shape(sl, bar_data,
                Inches(0.55), Inches(1.55),
                Inches(6.56), Inches(5.45),
                bar_colors=bar_cols,
                base_line=15, base_label="Azar ~15%")

# Panel derecho
card(sl, Inches(7.42), Inches(1.28), Inches(5.62), Inches(5.8), CARD)
lines_r = [
    ("5 NIVELES DE CERTEZA CIENTÍFICA", ACCENT, 12.5, True),
    ("", WHITE, 4, False),
    ("Son los 5 grados de EVIDENCIA que respaldan una alerta. Se diferencian de los niveles del mapa operacional.", SUBTITLE, 10, False),
    ("", WHITE, 8, False),
    ("N1 — FLAGRANCIA:   98.1%", RED,    11, True),
    ("Satelite confirma calor HOY + historial claro.", SUBTITLE, 9.5, False),
    ("", WHITE, 4, False),
    ("N2 — PATRON FUERTE: 86.4%", ORANGE, 11, True),
    ("Historia muy repetida, sin confirmacion satelital aun.", SUBTITLE, 9.5, False),
    ("", WHITE, 4, False),
    ("N3 — IA DETECTA:    72.3%", PURPLE, 11, True),
    ("La IA cruzo factores de riesgo: vias, pueblos, racha.", SUBTITLE, 9.5, False),
    ("", WHITE, 4, False),
    ("N4 — SEÑAL DEBIL:   44.7%", ACCENT, 11, True),
    ("Algun antecedente existe pero no es fuerte.", SUBTITLE, 9.5, False),
    ("", WHITE, 4, False),
    ("N5 — SIN DATOS:     11.9%", GREEN,  11, True),
    ("Sin historial previo. Linea base del sistema.", SUBTITLE, 9.5, False),
]
add_multiline_box(sl, lines_r,
                  Inches(7.6), Inches(1.4), Inches(5.28), Inches(5.5),
                  line_space_pct=128)

footer_line(sl)
print("  ✅ Slide 6 — Certeza Científica")

# =============================================================================
#  SLIDE 7A — LOS 4 NIVELES DE CRITICIDAD Y SUS RANGOS DE TIEMPO
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "LOS 4 NIVELES DE CRITICIDAD DE SATICA",
    "Cada hacienda tiene un ciclo histórico propio. SATICA mide cuán cerca está de ese ciclo en este momento.")

# Explicacion del concepto (panel superior izquierdo)
card(sl, Inches(0.48), Inches(1.28), Inches(5.6), Inches(2.15), CARD)
add_multiline_box(sl, [
    ("¿CÓMO FUNCIONA EL RELOJ?", ACCENT, 12, True),
    ("", WHITE, 5, False),
    ("Cada hacienda tiene su propio ciclo: el tiempo promedio entre un incendio y el siguiente, calculado con su historial real.", SUBTITLE, 10.5, False),
    ("", WHITE, 5, False),
    ("SATICA mide cuántos meses faltan o sobran respecto a ese ciclo y asigna el nivel de criticidad.", WHITE, 10.5, False),
], Inches(0.62), Inches(1.4), Inches(5.3), Inches(1.95), line_space_pct=130)

# Ejemplo concreto (panel inferior izquierdo)
card(sl, Inches(0.48), Inches(3.52), Inches(5.6), Inches(3.56), CARD)
add_multiline_box(sl, [
    ("EJEMPLO REAL", GREEN, 12, True),
    ("", WHITE, 5, False),
    ("Hacienda EL PARAÍSO arde en promedio cada 8 meses.", WHITE, 11, False),
    ("Último incendio: enero 2025.", WHITE, 11, False),
    ("Ciclo vence: septiembre 2025.", WHITE, 11, False),
    ("", WHITE, 6, False),
    ("En julio 2025:\nFaltan 2 meses → NIVEL ALTO", ORANGE, 11, True),
    ("En agosto 2025:\nFalta 1 mes → NIVEL CRÍTICO", RED, 11, True),
    ("En diciembre 2025:\nPasaron 3 meses sin arder → MITIGADO", GREEN, 11, True),
], Inches(0.62), Inches(3.65), Inches(5.3), Inches(3.3), line_space_pct=130)

# Los 4 niveles (panel derecho — columnas)
niveles_op = [
    (RED,    "CRÍTICO",         "#c0392b",
     "± 1 mes del ciclo",
     "La hacienda está en o muy cerca de su fecha histórica de incendio.",
     "Accion INMEDIATA · Visita urgente · Alerta Telegram activada"),
    (ORANGE, "ALTO",            "#e67e22",
     "1 a 2 meses antes/después",
     "Se acerca o acaba de pasar su ciclo. Alta probabilidad de evento inminente.",
     "Visita preventiva urgente · Ronda cortafuego · Verificacion NDVI"),
    (YELLOW, "OBSERVACIÓN",     "#f1c40f",
     "2 a 3 meses antes/después",
     "En el horizonte de riesgo. Se aproxima su temporada de mayor vulnerabilidad.",
     "Monitoreo activo · Incluir en ruta semanal de vigilancia"),
    (GREEN,  "BAJO / MITIGADO", "#27ae60",
     "Más de 3 meses",
     "Lejos de su ciclo o ya superó el período de riesgo sin que ardiera (éxito de gestión).",
     "Vigilancia pasiva · Registro de éxito ambiental si hubo visita previa"),
]

nv_w = Inches(1.73)
nv_x0 = Inches(6.28)
nv_gap = Inches(1.76)

for i, (col, nombre, hexcol, rango, desc, accion) in enumerate(niveles_op):
    lx = nv_x0 + i * nv_gap

    # Cabecera coloreada
    hd = sl.shapes.add_shape(1, lx, Inches(1.28), nv_w, Inches(0.85))
    hd.fill.solid(); hd.fill.fore_color.rgb = col; hd.line.fill.background()
    add_text_box(sl, nombre, lx+Inches(0.04), Inches(1.3),
                 nv_w-Inches(0.08), Inches(0.5),
                 font_size=12.5, color=BG, bold=True, align=PP_ALIGN.CENTER)

    # Rango de tiempo badge
    card(sl, lx, Inches(2.16), nv_w, Inches(0.52), RGBColor(0x16,0x27,0x3a))
    add_rect(sl, lx, Inches(2.16), Inches(0.08), Inches(0.52), col)
    add_text_box(sl, rango,
                 lx+Inches(0.14), Inches(2.2), nv_w-Inches(0.18), Inches(0.44),
                 font_size=8.5, color=col, bold=True, align=PP_ALIGN.CENTER)

    # Descripción
    card(sl, lx, Inches(2.72), nv_w, Inches(2.2), CARD)
    add_text_box(sl, desc,
                 lx+Inches(0.08), Inches(2.8), nv_w-Inches(0.16), Inches(2.0),
                 font_size=8.8, color=SUBTITLE, line_spacing=1.3)

    # Acción
    card(sl, lx, Inches(4.96), nv_w, Inches(2.12), RGBColor(0x16,0x27,0x3a))
    add_rect(sl, lx, Inches(4.96), Inches(0.08), Inches(2.12), col)
    add_text_box(sl, "ACCION:", lx+Inches(0.14), Inches(5.0),
                 nv_w-Inches(0.18), Inches(0.25),
                 font_size=7.5, color=FOOTER_C, bold=True)
    add_text_box(sl, accion, lx+Inches(0.14), Inches(5.25),
                 nv_w-Inches(0.18), Inches(1.75),
                 font_size=8.2, color=WHITE, line_spacing=1.3)

footer_line(sl)
print("  ✅ Slide 7A — 4 Niveles de Criticidad")



# =============================================================================
#  SLIDE 7B — DOS LENTES, UN SOLO SISTEMA
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "DOS LENTES — UN SOLO SISTEMA",
    "SATICA evalúa cada hacienda con dos criterios distintos que se complementan para emitir la alerta final")

# ── Centro: título de los dos lentes ────────────────────────────────────────
add_rect(sl, Inches(0.48), Inches(1.28), Inches(12.6), Inches(0.58), DARK_LINE)
add_text_box(sl, "LENTE 1 — RELOJ DEL CICLO",
             Inches(0.55), Inches(1.32), Inches(5.5), Inches(0.44),
             font_size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, "LENTE 2 — CERTEZA CIENTÍFICA",
             Inches(7.1), Inches(1.32), Inches(5.9), Inches(0.44),
             font_size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

# Subtítulos lentes
add_text_box(sl, "¿Cuándo actuar? (basado en tiempo del ciclo histórico)",
             Inches(0.55), Inches(1.78), Inches(5.8), Inches(0.25),
             font_size=8.5, color=SUBTITLE, align=PP_ALIGN.CENTER)
add_text_box(sl, "¿Con qué certeza? (basado en evidencia y modelo IA)",
             Inches(7.1), Inches(1.78), Inches(5.9), Inches(0.25),
             font_size=8.5, color=SUBTITLE, align=PP_ALIGN.CENTER)

# Separador central
add_rect(sl, Inches(6.5), Inches(1.28), Inches(0.08), Inches(5.8), DARK_LINE)

# ── LENTE 1: 4 niveles operacionales ────────────────────────────────────────
op_niveles = [
    (RED,    "CRÍTICO",     "± 1 mes del ciclo",       "Visita INMEDIATA ese día"),
    (ORANGE, "ALTO",        "1 a 2 meses antes/después","Visita preventiva urgente"),
    (YELLOW, "OBSERVACIÓN", "2 a 3 meses antes/después","Ruta de vigilancia semanal"),
    (GREEN,  "BAJO/MITIGADO","Más de 3 meses",          "Vigilancia pasiva"),
]
for j, (col, nombre, rango, accion) in enumerate(op_niveles):
    oy = Inches(2.1) + j * Inches(1.22)
    # Badge
    bd = sl.shapes.add_shape(1, Inches(0.55), oy, Inches(1.35), Inches(0.38))
    bd.fill.solid(); bd.fill.fore_color.rgb = col; bd.line.fill.background()
    add_text_box(sl, nombre, Inches(0.57), oy + Inches(0.03),
                 Inches(1.3), Inches(0.32), font_size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, rango,
                 Inches(2.0), oy + Inches(0.03), Inches(2.0), Inches(0.3),
                 font_size=9.5, color=col, bold=True)
    add_text_box(sl, accion,
                 Inches(2.0), oy + Inches(0.36), Inches(4.3), Inches(0.28),
                 font_size=8.5, color=SUBTITLE)
    add_rect(sl, Inches(0.55), oy + Inches(1.16), Inches(5.85), Pt(1), DARK_LINE)

# ── LENTE 2: 5 niveles científicos ──────────────────────────────────────────
sci_niveles = [
    (RED,    "N1 FLAGRANCIA",   "98.1%", "Satelite + historial fuerte"),
    (ORANGE, "N2 PATRON",       "86.4%", "Historial muy repetido"),
    (PURPLE, "N3 IA",           "72.3%", "Modelo IA detecta riesgo"),
    (ACCENT, "N4 SEÑAL DEBIL",  "44.7%", "Antecedente leve"),
    (GREEN,  "N5 SIN DATOS",    "11.9%", "Sin historial previo"),
]
for k, (col, nombre, pct, desc) in enumerate(sci_niveles):
    sy = Inches(2.1) + k * Inches(0.97)
    bd2 = sl.shapes.add_shape(1, Inches(6.65), sy, Inches(1.6), Inches(0.35))
    bd2.fill.solid(); bd2.fill.fore_color.rgb = col; bd2.line.fill.background()
    add_text_box(sl, nombre, Inches(6.67), sy + Inches(0.02),
                 Inches(1.55), Inches(0.3), font_size=9, color=BG, bold=True, align=PP_ALIGN.CENTER)
    # Barra precision
    track_w = Inches(3.5)
    add_rect(sl, Inches(8.35), sy + Inches(0.03), track_w, Inches(0.3), RGBColor(0x33,0x41,0x55))
    fw2 = track_w * (float(pct.replace("%","")) / 100)
    prg2 = sl.shapes.add_shape(1, Inches(8.35), sy + Inches(0.03), fw2, Inches(0.3))
    prg2.fill.solid(); prg2.fill.fore_color.rgb = col; prg2.line.fill.background()
    add_text_box(sl, pct, Inches(8.38), sy + Inches(0.04), fw2, Inches(0.25),
                 font_size=9, color=BG, bold=True)
    add_text_box(sl, desc, Inches(11.95), sy + Inches(0.03), Inches(1.1), Inches(0.32),
                 font_size=7.5, color=SUBTITLE)

# ── Banner inferior: cuando los dos lentes coinciden ─────────────────────────
add_rect(sl, Inches(0.48), Inches(6.52), Inches(12.6), Inches(0.56), RED)
add_text_box(sl,
    "CUANDO COINCIDEN: el predio es CRÍTICO en el reloj  +  N1 FLAGRANCIA en certeza  →  Alerta Telegram AUTOMÁTICA. El técnico debe ir ese día.",
    Inches(0.62), Inches(6.57), Inches(12.3), Inches(0.46),
    font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

footer_line(sl)
print("  ✅ Slide 7B — Dos lentes, un sistema")


# =============================================================================
#  SLIDE 8 — DASHBOARD: 5 MÓDULOS
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "EL DASHBOARD: 5 MÓDULOS OPERATIVOS",
    "Una herramienta táctica y estratégica para técnicos y directivos de la CVC")

modulos = [
    ("01\nRADAR\nPREVENTIVO", ACCENT,
     ["Mapa interactivo haciendas.", "Colores: CRÍTICO/ALTO/OBS/BAJO.",
      "Capas: GOES-16, HYSPLIT,\nRest. CVC.", "Pop-ups con historial y\nfecha estimada incendio."]),
    ("02\nRED DE\nRIESGO", PURPLE,
     ["Grafo Hacienda ↔ Ingenio.", "Concentración de riesgo\npor empresa.",
      "Identifica ingenios con\nmayor exposición.",
      "Gestión compromisos legales."]),
    ("03\nANALÍTICA\nHISTÓRICA", GREEN,
     ["Series de tiempo 2019–2026.", "Filtros: año, mes, municipio.",
      "Ranking Top 10 predios.", "Descarga PNG y CSV."]),
    ("04\nBASE DE\nDATOS", ORANGE,
     ["Tabla completa → Excel.", "Fechas, coordenadas, ciclos.",
      "Historial visitas técnicas.", "Anclas Operativas GPS."]),
    ("05\nSIN\nGEORREF.", RED,
     ["Gestión 'Punto Ciego'.", "Haciendas sin shapefile.",
      "Ingreso GPS manual.", "Integradas al análisis."]),
]

card_w2 = Inches(2.36)
for i, (title, col, bullets) in enumerate(modulos):
    lx = Inches(0.48) + i * Inches(2.57)
    card(sl, lx, Inches(1.3), card_w2, Inches(5.78), CARD)
    # Cabecera
    sh = slide.shapes.add_shape if False else sl.shapes.add_shape
    hd = sl.shapes.add_shape(1, lx, Inches(1.3), card_w2, Inches(1.5))
    hd.fill.solid(); hd.fill.fore_color.rgb = col; hd.line.fill.background()
    add_text_box(sl, title,
                 lx+Inches(0.05), Inches(1.32), card_w2-Inches(0.1), Inches(1.44),
                 font_size=11.5, color=BG, bold=True, align=PP_ALIGN.CENTER)
    lines_m = []
    for b in bullets:
        lines_m.append(("• " + b, SUBTITLE, 10, False))
        lines_m.append(("", WHITE, 3, False))
    add_multiline_box(sl, lines_m,
                      lx+Inches(0.1), Inches(2.88),
                      card_w2-Inches(0.2), Inches(4.1),
                      line_space_pct=130)

footer_line(sl)
print("  ✅ Slide 8 — Dashboard")

# =============================================================================
#  SLIDE 9 — AUTOMATIZACIÓN (El Centinela)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "AUTOMATIZACIÓN TOTAL: EL CENTINELA",
    "SATICA opera 24 / 7 sin intervención humana — sistema serverless en la nube de GitHub")

# Panel izquierdo
card(sl, Inches(0.48), Inches(1.28), Inches(6.0), Inches(5.8), CARD)
lines_auto = [
    ("⚙️  ROBOT EN LA NUBE (GITHUB ACTIONS)", GREEN, 12.5, True),
    ("", WHITE, 5, False),
    ("El archivo centinela.yml dispara un robot cada 15 minutos en servidores Ubuntu de Microsoft (GitHub). No requiere servidor propio ni costo de infraestructura.", SUBTITLE, 10.5, False),
    ("", WHITE, 8, False),
    ("• Consulta NASA FIRMS (VIIRS + MODIS real-time).", WHITE, 11, False),
    ("• Cruza focos contra mapa catastral de suertes.", WHITE, 11, False),
    ("• Calcula nivel de riesgo con modelo IA V9.", WHITE, 11, False),
    ("• Envía alerta Telegram si hay confirmación.", WHITE, 11, False),
    ("", WHITE, 8, False),
    ("🔒  Llaves API y tokens cifrados como GitHub Secrets.", FOOTER_C, 10, True),
]
add_multiline_box(sl, lines_auto,
                  Inches(0.65), Inches(1.42), Inches(5.68), Inches(5.5),
                  line_space_pct=135)

# Panel derecho — Ficha Telegram simulada
card(sl, Inches(6.74), Inches(1.28), Inches(6.3), Inches(5.8), CARD)
# Encabezado ficha
fh = sl.shapes.add_shape(1, Inches(6.74), Inches(1.28), Inches(6.3), Inches(0.7))
fh.fill.solid(); fh.fill.fore_color.rgb = RED; fh.line.fill.background()
add_text_box(sl, "🔥  ALERTA SATICA V2.0 · Certeza Factual 98.1%",
             Inches(6.85), Inches(1.32), Inches(6.1), Inches(0.56),
             font_size=11.5, color=WHITE, bold=True)

lines_tg = [
    ("Hacienda:       EL PARAISO", WHITE, 11.5, False),
    ("Suerte:             000123", WHITE, 11.5, False),
    ("Coordenadas:  3.8841, -76.4712", WHITE, 11.5, False),
    ("Estado:           Evento Térmico < 15 min", WHITE, 11.5, False),
    ("Riesgo:           CRÍTICO 🔴", RED, 12, True),
    ("", WHITE, 5, False),
    ("📡  SMART SYNC EN EL DASHBOARD", ACCENT, 11.5, True),
    ("", WHITE, 4, False),
    ("Al abrir el Dashboard, verifica si los datos tienen más de 30 min. Si están viejos, actualiza automáticamente desde NASA antes de mostrar el mapa.", SUBTITLE, 10.5, False),
    ("", WHITE, 6, False),
    ("🔗  Deep Link incluido → abre el mapa directamente centrado en el predio afectado.", GREEN, 10.5, True),
]
add_multiline_box(sl, lines_tg,
                  Inches(6.88), Inches(2.1), Inches(5.96), Inches(4.85),
                  line_space_pct=130)

footer_line(sl)
print("  ✅ Slide 9 — Automatización")

# =============================================================================
#  SLIDE 10 — BLINDAJE JURÍDICO
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "MÁS QUE UN MAPA: BLINDAJE JURÍDICO",
    "SATICA no solo alerta — documenta, justifica y da respaldo legal a cada acción de la CVC")

pilares = [
    ("Resolución 0741 de 2016", RED,
     "SATICA da sustento técnico a las visitas preventivas que exige esta norma. Cada alerta incluye coordenadas, fecha y nivel de confianza certificado científicamente."),
    ("Actas de Inspección", ORANGE,
     "El boletín ejecutivo (PDF/DOCX) por cada hacienda describe su historial, índice de vulnerabilidad y justifica técnico-científicamente la necesidad de la visita de inspección."),
    ("Métrica de Éxito Ambiental", GREEN,
     "Si una hacienda en nivel CRÍTICO recibe visita preventiva y supera su ciclo predicho sin arder, SATICA lo registra como 'Incendio Evitado': indicador de gestión ambiental de la CVC."),
    ("Exportación GIS Profesional", ACCENT,
     "Las capas KML exportadas llevan estilos de riesgo embebidos (rojo = Crítico, naranja = Alto). Se cargan directamente en QGIS o ArcGIS Pro sin configuración adicional."),
]

grid = [(Inches(0.48), Inches(1.3)), (Inches(6.92), Inches(1.3)),
        (Inches(0.48), Inches(4.35)), (Inches(6.92), Inches(4.35))]

for (lx, ty), (title, col, desc) in zip(grid, pilares):
    card(sl, lx, ty, Inches(6.2), Inches(2.75), CARD)
    # Barra lateral de color
    add_rect(sl, lx, ty, Inches(0.12), Inches(2.75), col)
    add_text_box(sl, title, lx+Inches(0.22), ty+Inches(0.1),
                 Inches(5.85), Inches(0.45), font_size=13.5, color=col, bold=True)
    add_text_box(sl, desc, lx+Inches(0.22), ty+Inches(0.6),
                 Inches(5.85), Inches(2.05), font_size=11, color=SUBTITLE, line_spacing=1.35)

footer_line(sl)
print("  ✅ Slide 10 — Blindaje Jurídico")

# =============================================================================
#  SLIDE 11 — ROADMAP
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)
accent_bar_left(sl)
title_slide_header(sl, "SIGUIENTES PASOS: SATICA V3.0+",
    "Los cimientos están listos. Estas expansiones multiplicarán el impacto del sistema.")

roadmap = [
    ("01", "Georref.\nCompleta",      GREEN,  "EN CURSO",
     "Completar polígonos de haciendas sin shapefile. Se usa módulo GPS manual activo."),
    ("02", "Portal Web\nCVC",         GREEN,  "EN CURSO",
     "Despliegue en servidor con acceso autenticado para todos los DAR de la CVC."),
    ("03", "App Móvil\nde Campo",     ACCENT, "PLANIFICADO",
     "PWA para técnicos: alertas, registro de visitas y evidencia fotográfica en terreno."),
    ("04", "Integración\nCDIAC",      ACCENT, "PLANIFICADO",
     "Datos climáticos IDEAM: sequía, ETP, déficit hídrico como features del modelo IA."),
    ("05", "Régimen de\nQuemas",      PURPLE, "PLANIFICADO",
     "Módulo legal para solicitudes de quemas: evalúa restricciones CVC y viento actual."),
    ("06", "IA V10\nAutolearn",       PURPLE, "FUTURO",
     "Reentrenamiento automático del XGBoost cada trimestre con nuevos datos confirmados."),
]

# Línea timeline
add_rect(sl, Inches(0.8), Inches(2.55), Inches(12.1), Pt(3), DARK_LINE)

step_x = Inches(12.1) / 5

for i, (num, title, col, estado, desc) in enumerate(roadmap):
    cx = Inches(0.8) + i * step_x + step_x/2 - Inches(0.07)

    # Círculo en timeline
    cr = Inches(0.25)
    circ = sl.shapes.add_shape(9, cx-cr, Inches(2.42), cr*2, cr*2)
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    add_text_box(sl, num, cx-cr, Inches(2.48), cr*2, Inches(0.4),
                 font_size=9, color=BG, bold=True, align=PP_ALIGN.CENTER)

    # Tarjeta
    cx_c = cx - Inches(0.95)
    cy_c = Inches(1.25) if i % 2 == 0 else Inches(2.95)
    card(sl, cx_c, cy_c, Inches(1.9), Inches(1.12) if i%2==0 else Inches(3.6), CARD)
    add_text_box(sl, title, cx_c+Inches(0.08), cy_c+Inches(0.06),
                 Inches(1.74), Inches(0.48),
                 font_size=10.5, color=col, bold=True, align=PP_ALIGN.CENTER)
    # Estado badge
    col_est = GREEN if estado=="EN CURSO" else (ACCENT if estado=="PLANIFICADO" else PURPLE)
    add_text_box(sl, estado, cx_c+Inches(0.08), cy_c+Inches(0.52),
                 Inches(1.74), Inches(0.25),
                 font_size=7.5, color=col_est, bold=True, align=PP_ALIGN.CENTER)
    if i % 2 == 0:
        # Línea conector hacia abajo
        add_rect(sl, cx, cy_c+Inches(1.12), Pt(2), Inches(0.18), DARK_LINE)

    # Descripción debajo del timeline
    add_text_box(sl, desc,
                 cx_c, Inches(3.25), Inches(1.9), Inches(2.9),
                 font_size=8.5, color=SUBTITLE, align=PP_ALIGN.CENTER, line_spacing=1.25)

footer_line(sl)
print("  ✅ Slide 11 — Roadmap")

# =============================================================================
#  SLIDE 12 — CIERRE
# =============================================================================
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, BG)

# Gradiente de fondo diagonal
gradient_rect(sl, 0, 0, W, H,
              RGBColor(0x0c,0x12,0x22), RGBColor(0x0f,0x28,0x42), angle=135)

# Franjas superior e inferior
for ty in [0, H - Inches(0.38)]:
    gradient_rect(sl, 0, ty, W, Inches(0.38),
                  ACCENT, RGBColor(0x06,0x82,0xa8), angle=0)

# Texto central
add_text_box(sl, "SATICA",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.7),
             font_size=96, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(sl, "protege el Valle del Cauca",
             Inches(0.5), Inches(2.72), Inches(12.3), Inches(0.65),
             font_size=26, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(3.0), Inches(3.45), Inches(7.33), Pt(2), ACCENT)

add_text_box(sl,
    "«Con la tecnología correcta, un técnico de la CVC tiene hoy más poder de\ndetección que una brigada completa de hace 10 años.»",
    Inches(1.2), Inches(3.58), Inches(10.9), Inches(0.85),
    font_size=12.5, color=SUBTITLE, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.4)

# 3 stats finales
stats_c = [
    ("5 satélites\nintegrados", ACCENT),
    ("98.1%\ncerteza factual", RED),
    ("15 minutos\nde respuesta", GREEN),
]
for i, (lbl, col) in enumerate(stats_c):
    sx = Inches(1.5) + i * Inches(3.6)
    card(sl, sx, Inches(4.65), Inches(3.2), Inches(1.45), CARD)
    add_text_box(sl, lbl, sx+Inches(0.1), Inches(4.75),
                 Inches(3.0), Inches(1.25),
                 font_size=16, color=col, bold=True, align=PP_ALIGN.CENTER)

add_text_box(sl,
    "Sistema de Alertas Tempranas de Incendios en Caña de Azúcar  ·  CVC — DAR Suroriente  ·  Abril 2026",
    Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.4),
    font_size=9, color=FOOTER_C, align=PP_ALIGN.CENTER)

print("  ✅ Slide 12 — Cierre")

# =============================================================================
#  GUARDAR
# =============================================================================
output = "PRESENTACION_SATICA_V2.pptx"
prs.save(output)

print()
print("╔══════════════════════════════════════════════════════════════╗")
print("║  ✅  PRESENTACIÓN GENERADA EXITOSAMENTE                     ║")
print(f"║  📄  Archivo: {output:<46}║")
print("║  📊  Diapositivas: 12  ·  Diseño: Dark Premium              ║")
print("╚══════════════════════════════════════════════════════════════╝")
print(f"\nAbra: {output}")
